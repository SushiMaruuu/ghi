import streamlit as st
from pptx import Presentation
import openai
from dotenv import load_dotenv
import os

# Load environment variables from .env file
load_dotenv()

# Set your OpenAI API key from the .env file
openai.api_key = os.getenv("OPENAI_API_KEY")

# Function to clean text (replace unsupported characters, e.g., en-dash)
def clean_text(text):
    # Remove unwanted characters and filter out very short lines
    unwanted_phrases = ['Slide', 'OCTOBER', 'Short URL']
    for phrase in unwanted_phrases:
        text = text.replace(phrase, '')
    text = text.replace('\u2013', '-')  # Replace en-dash with a simple dash
    return text

# Function to generate multiple-choice questions and summary using OpenAI GPT model
def generate_mcqs_and_summary(text):
    # Limit text input size for OpenAI (truncate if too long)
    truncated_text = text[:4000]  # OpenAI's models have a token limit, so we truncate

    # Call OpenAI's GPT model to generate MCQs and a summary
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",  # Use GPT-3.5 Turbo
            messages=[
                {"role": "system", "content": "You are a helpful assistant that generates multiple-choice questions with answers and a summary from content."},
                {"role": "user", "content": f"Generate exactly 10 multiple-choice questions with four options each (one correct answer), and provide the correct answer after each question. Then, provide a separate summary of the following content. Separate the questions and summary clearly:\n\n{truncated_text}\n\n--- Summary ---"}
            ],
            max_tokens=1000,  # Allow more space for questions and summary
            temperature=0.7,
            n=1,
        )
        
        # Extract the generated text from the response
        generated_text = response['choices'][0]['message']['content'].strip()
        
        # Split the response into questions and summary using the separator
        parts = generated_text.split("--- Summary ---")
        questions_part = parts[0].strip()
        summary = parts[1].strip() if len(parts) > 1 else "No summary generated."

        # Ensure questions are in the correct format
        mcqs = questions_part.split("\n\n")  # Each question/answer block separated by double newline

        return mcqs, summary

    except Exception as e:
        st.error(f"Error generating multiple-choice questions and summary: {str(e)}")
        return [], ""

# Streamlit app
st.title("PowerPoint Text Extractor and MCQ Generator")

# File uploader for PPTX files only
uploaded_file = st.file_uploader("Upload a PowerPoint (PPTX)", type=["pptx"])

if uploaded_file is not None:
    st.write(f"Filename: {uploaded_file.name}")
    
    # Load the PowerPoint file
    presentation = Presentation(uploaded_file)
    all_text = []

    # Extract text from each slide
    for i, slide in enumerate(presentation.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):  # Only extract text-containing shapes
                slide_text.append(shape.text)
        
        # Clean and join the text for each slide
        cleaned_slide_text = clean_text("\n".join(slide_text))
        all_text.append(f"*Slide {i + 1}:*\n{cleaned_slide_text}\n")

    # Display the extracted text from all slides
    st.subheader("Extracted Text from PowerPoint")
    for slide_text in all_text:
        st.write(slide_text)

    # Optionally, show a message if no text is found in the slides
    if not any(slide_text for slide_text in all_text):
        st.write("No text found in the PowerPoint.")
    
    # Concatenate all slide texts for generating questions
    full_text = " ".join([clean_text(slide) for slide in all_text])

    # Button to generate multiple-choice questions and summary
    if st.button("Generate MCQs and Summary"):
        mcqs, summary = generate_mcqs_and_summary(full_text)

        st.subheader("Generated Multiple-Choice Questions with Answers")
        for mcq in mcqs:
            st.write(mcq)

        st.subheader("Generated Summary")
        if summary != "No summary generated.":
            st.write(summary)
